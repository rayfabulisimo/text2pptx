import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
import requests  # For fetching cartoon images (or use OpenAI API)

# --- CONFIGURATION ---
SLIDE_SIZE = (13.33, 7.5)  # 16:9 inches
THEME_COLOR = RGBColor(100, 100, 100)  # Slate gray (for sidebar)
PASTEL_BACKGROUND = RGBColor(249, 249, 249)  # Cream background
FONT_TITLE = "Arial Black"  # Bold hierarchy
FONT_BODY = "Arial"         # Light body text

# --- HELPER FUNCTIONS ---
def fetch_cartoon_image(prompt: str, output_dir: str = "cartoon_images"):
    """Fetch a cartoon image from DALL-E/Midjourney (or use a placeholder)."""
    os.makedirs(output_dir, exist_ok=True)
    image_path = f"{output_dir}/{hash(prompt)}.png"

    if not os.path.exists(image_path):
        # Create a simple placeholder image using PIL
        from PIL import Image, ImageDraw, ImageFont
        img = Image.new('RGB', (800, 600), color=(249, 249, 249))
        d = ImageDraw.Draw(img)
        d.text((400, 300), "Placeholder Image", fill=(100, 100, 100), anchor="mm")
        img.save(image_path)
        print(f"[DEBUG] Created placeholder image: {image_path}")

    return image_path

def create_slide_layout(slide, layout_type: str, slide_json):
    """Apply the correct layout (SPLIT_SCREEN, CENTER_CARD, SIDEBAR)."""
    if layout_type == "SPLIT_SCREEN":
        return split_screen_layout(slide, slide_json)
    elif layout_type == "CENTER_CARD":
        return center_card_layout(slide, slide_json)
    elif layout_type == "SIDEBAR":
        return sidebar_layout(slide, slide_json)
    else:
        raise ValueError(f"Unknown layout: {layout_type}")

def split_screen_layout(slide, slide_json):
    """Left: Text, Right: Image (with optional sidebar)."""
    # Text Box (Left)
    left = Inches(slide_json["content"]["text_box"]["left"])
    top = Inches(slide_json["content"]["text_box"]["top"])
    width = Inches(slide_json["content"]["text_box"]["width"])
    height = Inches(3.5)  # Fixed height for text

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = slide_json["content"]["title"]
    tf.word_wrap = True
    tf.paragraphs[0].font.name = FONT_TITLE
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True

    # Add bullets
    for bullet in slide_json["content"]["bullets"]:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(14)

    # Image (Right)
    image_path = fetch_cartoon_image(slide_json["visual"]["image_prompt"])
    img_left = Inches(slide_json["visual"]["position"]["left"])
    img_top = Inches(slide_json["visual"]["position"]["top"])
    img_width = Inches(slide_json["visual"]["position"]["width"])

    slide.shapes.add_picture(image_path, img_left, img_top, img_width, Inches(3.5))

def center_card_layout(slide, slide_json):
    """Centered text with a cartoon below."""
    # Title (Centered)
    title_left = Inches((SLIDE_SIZE[0] - 6) / 2)  # Centered
    title_top = Inches(1.5)
    title_width = Inches(6)

    txBox = slide.shapes.add_textbox(title_left, title_top, title_width, Inches(2))
    tf = txBox.text_frame
    tf.text = slide_json["content"]["title"]
    tf.word_wrap = True
    tf.paragraphs[0].font.name = FONT_TITLE
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True

    # Cartoon (Below)
    image_path = fetch_cartoon_image(slide_json["visual"]["image_prompt"])
    img_top = Inches(4.5)  # Below title
    img_left = Inches((SLIDE_SIZE[0] - 6) / 2)  # Centered
    img_width = Inches(6)

    slide.shapes.add_picture(image_path, img_left, img_top, img_width, Inches(3.5))

def sidebar_layout(slide, slide_json):
    """Sidebar (left) + Content (right)."""
    # Sidebar (Left)
    sidebar_width = Inches(2)
    sidebar_left = Inches(0)
    sidebar_top = Inches(0)
    sidebar_height = Inches(SLIDE_SIZE[1])

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        sidebar_left, sidebar_top,
        sidebar_width, sidebar_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME_COLOR

    # Content (Right)
    content_left = Inches(2.5)
    content_top = Inches(1.5)
    content_width = SLIDE_SIZE[0] - Inches(3.5)

    txBox = slide.shapes.add_textbox(content_left, content_top, content_width, Inches(5))
    tf = txBox.text_frame
    tf.text = slide_json["content"]["title"]
    tf.word_wrap = True
    tf.paragraphs[0].font.name = FONT_TITLE
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True

    # Add bullets
    for bullet in slide_json["content"]["bullets"]:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(14)

    # Cartoon (Right side, below text)
    image_path = fetch_cartoon_image(slide_json["visual"]["image_prompt"])
    img_left = Inches(7.5)  # Right side
    img_top = Inches(4.5)
    img_width = Inches(5)

    slide.shapes.add_picture(image_path, img_left, img_top, img_width, Inches(3.5))

# --- MAIN FUNCTION ---
def generate_slides(json_path: str, output_pptx: str = "presentation.pptx"):
    """Generate slides from JSON design."""
    prs = Presentation()
    slide_layouts = prs.slide_layouts

    with open(json_path, "r") as f:
        slides_data = json.load(f)

    for slide_data in slides_data:
        slide = prs.slides.add_slide(slide_layouts[6])  # Blank slide
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = PASTEL_BACKGROUND

        # Apply layout
        create_slide_layout(slide, slide_data["layout_type"], slide_data)

    prs.save(output_pptx)
    print(f"Presentation saved to {output_pptx}")

# --- EXAMPLE USAGE ---
if __name__ == "__main__":
    # Example JSON (replace with your actual JSON)
    example_json = """
    [
        {
            "slide_number": 1,
            "layout_type": "SPLIT_SCREEN",
            "background_color": "#f9f9f9",
            "content": {
                "title": "Introduction to NotebookLM",
                "bullets": ["Minimalist design", "High-contrast visuals"],
                "text_box": {"top": 2.0, "left": 1.0, "width": 5.0}
            },
            "visual": {
                "image_prompt": "Minimalist 2D flat-vector illustration of a notebook with a lightbulb, pastel cream and slate palette, professional corporate-chic style, isolated on a solid background.",
                "position": {"top": 1.5, "left": 7.0, "width": 5.0}
            }
        }
    ]
    """

    # Save example JSON to a file
    with open("slides.json", "w") as f:
        f.write(example_json)

    # Generate slides
    generate_slides("slides.json")
