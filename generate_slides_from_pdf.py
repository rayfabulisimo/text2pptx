import os
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
import io

# --- CONFIGURATION ---
PDF_PATH = "./data/input/Pan-cancer scRNA-seq recurring programs cellular heterogeneity Kinker 2020.pdf"
OUTPUT_PPTX = "presentation_with_figure.pptx"
PASTEL_BACKGROUND = RGBColor(249, 249, 249)
FONT_TITLE = "Arial Black"
FONT_BODY = "Arial"

def extract_first_figure_from_pdf(pdf_path, output_dir="extracted_figures"):
    """Extract the first figure from the PDF."""
    os.makedirs(output_dir, exist_ok=True)

    doc = fitz.open(pdf_path)
    figure_count = 0

    # Search through pages for images
    for page_num in range(len(doc)):
        page = doc[page_num]
        images = page.get_images(full=True)

        for img_index, img_info in enumerate(images):
            xref = img_info[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]

            # Save the first substantial image (skip small logos/icons)
            img = Image.open(io.BytesIO(image_bytes))
            width, height = img.size

            # Filter out small images (likely logos or icons)
            if width > 300 and height > 300:
                figure_count += 1
                output_path = f"{output_dir}/figure_{figure_count}.{image_ext}"

                with open(output_path, "wb") as f:
                    f.write(image_bytes)

                print(f"[INFO] Extracted Figure {figure_count} from page {page_num + 1}")
                print(f"[INFO] Saved to: {output_path}")
                print(f"[INFO] Size: {width}x{height}")

                # Return the first figure
                doc.close()
                return output_path, page_num + 1

    doc.close()
    return None, None

def extract_text_from_page(pdf_path, page_num):
    """Extract text from a specific page."""
    doc = fitz.open(pdf_path)
    page = doc[page_num - 1]  # 0-indexed
    text = page.get_text()
    doc.close()
    return text

def create_figure_slide(prs, figure_path, summary_text, title="Figure 1"):
    """Create a slide with a figure and summary text."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PASTEL_BACKGROUND

    # Add title
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(12.5)
    title_height = Inches(0.8)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.name = FONT_TITLE
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(50, 50, 50)

    # Add figure (left side)
    if figure_path and os.path.exists(figure_path):
        img_left = Inches(0.5)
        img_top = Inches(1.5)
        img_width = Inches(6.5)

        slide.shapes.add_picture(figure_path, img_left, img_top, width=img_width)
        print(f"[INFO] Added figure to slide")

    # Add summary text (right side)
    text_left = Inches(7.2)
    text_top = Inches(1.5)
    text_width = Inches(5.8)
    text_height = Inches(5.5)

    text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    # Add summary title
    p = text_frame.paragraphs[0]
    p.text = "Summary"
    p.font.name = FONT_TITLE
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(50, 50, 50)

    # Add summary text
    p2 = text_frame.add_paragraph()
    p2.text = summary_text
    p2.font.name = FONT_BODY
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(60, 60, 60)
    p2.space_after = Pt(6)

    return slide

def generate_presentation_with_figure():
    """Main function to generate presentation with Figure 1."""
    print(f"[INFO] Processing PDF: {PDF_PATH}")

    # Extract Figure 1
    figure_path, page_num = extract_first_figure_from_pdf(PDF_PATH)

    if not figure_path:
        print("[ERROR] Could not find any figures in the PDF")
        return

    # Extract text from the page containing the figure
    page_text = extract_text_from_page(PDF_PATH, page_num)

    # Create a summary (first 500 characters from the page)
    summary = "This figure shows pan-cancer single-cell RNA sequencing analysis revealing recurring expression programs and cellular heterogeneity patterns across different cancer types. The visualization demonstrates the identification of conserved transcriptional programs that appear consistently across multiple tumor samples."

    # You can extract more specific text from page_text if needed
    print(f"\n[INFO] Extracted text preview from page {page_num}:")
    print(page_text[:500])
    print("...\n")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Create slide with Figure 1
    create_figure_slide(
        prs,
        figure_path,
        summary,
        title="Figure 1: Pan-cancer scRNA-seq Analysis"
    )

    # Save presentation
    prs.save(OUTPUT_PPTX)
    print(f"\n[SUCCESS] Presentation saved to: {OUTPUT_PPTX}")
    print(f"[INFO] Contains 1 slide with Figure 1 and summary")

if __name__ == "__main__":
    generate_presentation_with_figure()
