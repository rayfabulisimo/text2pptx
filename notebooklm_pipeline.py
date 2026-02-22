#!/usr/bin/env python3
"""
NotebookLM-Style Presentation Pipeline
1. Extract text from PDF pages
2. Send to Claude API to generate JSON layout
3. Build PowerPoint from JSON
"""

import json
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# System prompt for Claude API
SYSTEM_PROMPT = """Role: You are an expert Presentation Designer specializing in the "NotebookLM" aesthetic: minimalist, clean, high-contrast, and visually balanced.

Task: Analyze the provided PDF text and generate a JSON structure for slides.

Design Constraints:
* Layout Variety: Choose from: SPLIT_SCREEN (Text Left, Image Right), CENTER_CARD (Focused text in middle), or SIDEBAR (Accent color bar on left).
* Typography: Clear hierarchy. Title (Size 36, Bold), Bullets (Size 14).
* Coordinates: Provide placement in "Inches" for a 16:9 slide (13.33" x 7.5").
* Each slide has 3-5 bullet points, with 6-8 words each.

Output Format (JSON Array):
[
  {
    "slide_number": 1,
    "layout_type": "SPLIT_SCREEN",
    "background_color": "#f9f9f9",
    "accent_color": "#6495ED",
    "content": {
      "title": "Slide Title",
      "bullets": ["Point 1 with 6-8 words", "Point 2"],
      "text_box": {"top": 2.0, "left": 1.0, "width": 5.5, "height": 4.0}
    },
    "visual": {
      "image_prompt": "Minimalist 2D flat-vector illustration of [subject]",
      "position": {"top": 1.5, "left": 7.5, "width": 5.0, "height": 5.0}
    }
  }
]

Return ONLY the JSON array, no other text."""


def hex_to_rgb(hex_color):
    """Convert hex to RGBColor."""
    hex_color = hex_color.lstrip('#')
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )


def add_sidebar(slide, color, width=0.4):
    """Add colored sidebar."""
    sidebar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(width), Inches(7.5)
    )
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = hex_to_rgb(color)
    sidebar.line.fill.background()
    # Send to back
    slide.shapes._spTree.remove(sidebar._element)
    slide.shapes._spTree.insert(2, sidebar._element)


def create_slide_from_json(prs, slide_json):
    """Create slide from JSON specification."""
    layout_type = slide_json.get('layout_type', 'SPLIT_SCREEN')
    bg_color = slide_json.get('background_color', '#f9f9f9')

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(bg_color)

    # Add sidebar if needed
    if layout_type == 'SIDEBAR':
        accent_color = slide_json.get('accent_color', '#6495ED')
        add_sidebar(slide, accent_color)

    # Content
    content = slide_json.get('content', {})
    tb_spec = content.get('text_box', {})
    top = Inches(tb_spec.get('top', 2.0))
    left = Inches(tb_spec.get('left', 1.0))
    width = Inches(tb_spec.get('width', 5.5))
    height = Inches(tb_spec.get('height', 4.0))

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    # Title
    title = content.get('title', '')
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.name = "Arial"
        p.font.color.rgb = RGBColor(33, 33, 33)
        p.space_after = Pt(18)

    # Bullets
    bullets = content.get('bullets', [])
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(14)
        p.font.name = "Arial"
        p.font.color.rgb = RGBColor(33, 33, 33)
        p.space_before = Pt(10)
        p.line_spacing = 1.3

    # Visual placeholder
    visual = slide_json.get('visual')
    if visual:
        pos = visual.get('position', {})
        v_top = Inches(pos.get('top', 1.5))
        v_left = Inches(pos.get('left', 7.5))
        v_width = Inches(pos.get('width', 5.0))
        v_height = Inches(pos.get('height', 5.0))

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            v_left, v_top, v_width, v_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(230, 235, 245)
        shape.line.color.rgb = RGBColor(180, 190, 210)
        shape.line.width = Pt(1)

        # Add visual indicator
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        prompt_text = visual.get('image_prompt', 'Visual')[:50]
        p.text = f"[Visual]\n{prompt_text}..."
        p.font.size = Pt(10)
        p.font.italic = True
        p.font.color.rgb = RGBColor(100, 110, 130)
        p.alignment = PP_ALIGN.CENTER


def generate_json_with_mock_llm(pdf_text):
    """
    Mock LLM response - generates sample JSON.
    In production, this would call Claude API.
    """
    # For testing, create a simple JSON response
    slides_json = [
        {
            "slide_number": 1,
            "layout_type": "SPLIT_SCREEN",
            "background_color": "#f9f9f9",
            "accent_color": "#6495ED",
            "content": {
                "title": "Pan-Cancer Single-Cell Analysis",
                "bullets": [
                    "198 cancer cell lines profiled comprehensively",
                    "53,513 individual cells analyzed at resolution",
                    "22 distinct cancer types represented"
                ],
                "text_box": {"top": 2.0, "left": 1.0, "width": 5.5, "height": 4.0}
            },
            "visual": {
                "image_prompt": "Minimalist 2D flat-vector illustration of laboratory microscope and cell culture dishes, pastel cream and slate palette, professional corporate-chic style",
                "position": {"top": 1.5, "left": 7.5, "width": 5.0, "height": 5.0}
            }
        },
        {
            "slide_number": 2,
            "layout_type": "SIDEBAR",
            "background_color": "#ffffff",
            "accent_color": "#6495ED",
            "content": {
                "title": "12 Recurrent Heterogeneous Programs",
                "bullets": [
                    "Cell cycle G1/S and G2/M phases",
                    "Three distinct EMT program variants identified",
                    "Stress response and interferon pathways active",
                    "Programs mirror tumor heterogeneity patterns"
                ],
                "text_box": {"top": 2.5, "left": 1.5, "width": 10.0, "height": 4.0}
            }
        }
    ]

    return slides_json


def create_presentation_from_json(slides_json, output_file):
    """Create presentation from JSON."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    print(f"\nCreating NotebookLM-style presentation...")
    print(f"Slides: {len(slides_json)}")

    for idx, slide_json in enumerate(slides_json, 1):
        layout = slide_json.get('layout_type', 'UNKNOWN')
        title = slide_json.get('content', {}).get('title', 'Untitled')
        print(f"  Slide {idx}: {layout} - {title[:50]}...")
        create_slide_from_json(prs, slide_json)

    prs.save(output_file)
    print(f"\n✓ Saved: {output_file}")


def main():
    """Main pipeline."""
    # Read sample text from PDF
    pdf_file = "paper_text_full.txt"

    if Path(pdf_file).exists():
        with open(pdf_file, 'r') as f:
            # Read first 100 lines (about 2 pages)
            lines = f.readlines()[:100]
            pdf_text = ''.join(lines)
        print(f"Loaded {len(lines)} lines from {pdf_file}")
    else:
        pdf_text = "Sample cancer research text about cellular heterogeneity."
        print("Using mock text")

    # Generate JSON (mock - in production would call Claude API)
    print("\nGenerating slide layouts...")
    slides_json = generate_json_with_mock_llm(pdf_text)

    # Save JSON for inspection
    json_file = "slide_layouts.json"
    with open(json_file, 'w') as f:
        json.dump(slides_json, f, indent=2)
    print(f"Saved layout JSON to: {json_file}")

    # Create presentation
    output_file = "notebooklm_test_presentation.pptx"
    create_presentation_from_json(slides_json, output_file)


if __name__ == "__main__":
    main()
