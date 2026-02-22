#!/usr/bin/env python3
"""
Modern PPTX Generator with Beautiful Themes
Creates professional presentations with modern design
"""

import sys
import json
from pathlib import Path
from typing import Dict, List
import io

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Error: python-pptx not installed. Install with: pip install python-pptx")
    sys.exit(1)

try:
    from PIL import Image
except ImportError:
    print("Error: Pillow not installed. Install with: pip install Pillow")
    sys.exit(1)


# Modern color themes
THEMES = {
    "ocean": {
        "name": "Ocean Blue",
        "primary": RGBColor(13, 71, 161),      # Deep blue
        "secondary": RGBColor(25, 118, 210),   # Medium blue
        "accent": RGBColor(66, 165, 245),      # Light blue
        "background": RGBColor(245, 248, 250), # Very light blue-gray
        "text": RGBColor(33, 33, 33)           # Dark gray
    },
    "forest": {
        "name": "Forest Green",
        "primary": RGBColor(27, 94, 32),       # Deep green
        "secondary": RGBColor(56, 142, 60),    # Medium green
        "accent": RGBColor(129, 199, 132),     # Light green
        "background": RGBColor(245, 248, 245), # Very light green
        "text": RGBColor(33, 33, 33)
    },
    "sunset": {
        "name": "Sunset Orange",
        "primary": RGBColor(191, 54, 12),      # Deep orange
        "secondary": RGBColor(230, 74, 25),    # Medium orange
        "accent": RGBColor(255, 138, 101),     # Light orange
        "background": RGBColor(255, 248, 245), # Very light orange
        "text": RGBColor(33, 33, 33)
    },
    "purple": {
        "name": "Royal Purple",
        "primary": RGBColor(74, 20, 140),      # Deep purple
        "secondary": RGBColor(123, 31, 162),   # Medium purple
        "accent": RGBColor(186, 104, 200),     # Light purple
        "background": RGBColor(248, 245, 250), # Very light purple
        "text": RGBColor(33, 33, 33)
    },
    "slate": {
        "name": "Modern Slate",
        "primary": RGBColor(38, 50, 56),       # Dark slate
        "secondary": RGBColor(69, 90, 100),    # Medium slate
        "accent": RGBColor(96, 125, 139),      # Light slate
        "background": RGBColor(250, 250, 250), # Off-white
        "text": RGBColor(33, 33, 33)
    },
    "coral": {
        "name": "Coral Pink",
        "primary": RGBColor(194, 24, 91),      # Deep coral
        "secondary": RGBColor(236, 64, 122),   # Medium coral
        "accent": RGBColor(255, 138, 171),     # Light coral
        "background": RGBColor(255, 245, 248), # Very light pink
        "text": RGBColor(33, 33, 33)
    }
}


def add_background(slide, theme):
    """Add background - either image or gradient with proper styling."""
    # Check if theme has background image
    if "background_image" in theme and theme["background_image"]:
        try:
            from pathlib import Path
            img_path = Path(theme["background_image"])
            if img_path.exists():
                # Add background image
                left = Inches(0)
                top = Inches(0)
                height = Inches(7.5)
                width = Inches(10)

                picture = slide.shapes.add_picture(
                    str(img_path), left, top,
                    width=width, height=height
                )
                # Send to back
                slide.shapes._spTree.remove(picture._element)
                slide.shapes._spTree.insert(2, picture._element)
                return
        except Exception as e:
            print(f"  Warning: Could not add background image: {e}")

    # Rich gradient background
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 135  # Diagonal for more visual interest

    # Two-stop gradient with theme colors
    bg_color = theme.get("background", RGBColor(245, 248, 250))

    fill.gradient_stops[0].color.rgb = bg_color
    fill.gradient_stops[1].color.rgb = RGBColor(255, 255, 255)


def add_gradient_background(slide, theme):
    """Deprecated - use add_background instead."""
    add_background(slide, theme)


def create_title_slide(prs, title, subtitle, theme):
    """Create a modern title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Add background (gradient or image)
    add_background(slide, theme)

    # Add decorative shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(2),
        Inches(9), Inches(3.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme["primary"]
    shape.line.fill.background()
    shape.shadow.inherit = False

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.3),
        Inches(8), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)  # Larger, more impactful
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4),
        Inches(8), Inches(1)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)  # Slightly larger
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER


def create_content_slide(prs, title, content, image_data, theme):
    """Create a modern content slide with text and optional image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Add background (gradient or image)
    add_background(slide, theme)

    # Add header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = theme["primary"]
    header.line.fill.background()

    # Title in header
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        Inches(9), Inches(0.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)  # Larger, more readable
    p.font.bold = True
    p.font.name = "Calibri"  # Clean, professional font
    p.font.color.rgb = RGBColor(255, 255, 255)

    if image_data:
        # Layout with image on right
        text_width = 4.5
        text_left = 0.5

        img_width = 4.5
        img_left = 5.3
    else:
        # Full width text
        text_width = 9
        text_left = 0.5

    # Content text
    content_box = slide.shapes.add_textbox(
        Inches(text_left), Inches(1.2),
        Inches(text_width), Inches(5.8)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    # Parse and add content
    for line in content.split('\n'):
        line = line.strip()
        if not line:
            continue

        p = tf.add_paragraph() if len(tf.paragraphs) > 1 or tf.paragraphs[0].text else tf.paragraphs[0]

        # Check if it's a bullet point
        if line.startswith('- ') or line.startswith('• '):
            p.text = line[2:]
            p.level = 0
            p.font.size = Pt(16)  # Larger for better readability
        elif line.startswith('  - ') or line.startswith('  • '):
            p.text = line[4:]
            p.level = 1
            p.font.size = Pt(14)  # Sub-bullets slightly smaller
        else:
            p.text = line
            p.font.size = Pt(16)  # Larger body text

        p.font.name = "Calibri"  # Professional, clean font
        p.font.color.rgb = theme["text"]
        p.space_before = Pt(8)  # More spacing
        p.line_spacing = 1.2  # Better line spacing

    # Add image if provided
    if image_data:
        try:
            img_stream = io.BytesIO(image_data)
            img = Image.open(img_stream)
            orig_width, orig_height = img.size

            # More generous sizing - use full available space
            max_width = img_width  # 4.5 inches available
            max_height = 5.8  # Full height available

            aspect_ratio = orig_width / orig_height

            # Start with max width, then adjust
            new_width = max_width
            new_height = new_width / aspect_ratio

            # If too tall, scale down to fit height
            if new_height > max_height:
                new_height = max_height
                new_width = new_height * aspect_ratio

            # If still doesn't fit, ensure it does
            if new_width > max_width:
                new_width = max_width
                new_height = new_width / aspect_ratio

            # Center image vertically and horizontally in available space
            img_top = 1.2 + (max_height - new_height) / 2
            actual_left = img_left + (img_width - new_width) / 2

            img_stream.seek(0)
            pic = slide.shapes.add_picture(
                img_stream,
                Inches(actual_left), Inches(img_top),
                width=Inches(new_width),
                height=Inches(new_height)
            )

            print(f"    Added image: {new_width:.2f}\" x {new_height:.2f}\" (original: {orig_width}x{orig_height}px)")
        except Exception as e:
            print(f"  Warning: Could not add image: {e}")


def create_modern_presentation(
    summary_data: Dict,
    images: List[Dict],
    output_file: str,
    theme_name: str = "ocean",
    title: str = "Document Summary",
    custom_theme: Dict = None
):
    """
    Create a modern PowerPoint presentation.

    Args:
        summary_data: Dict with 'sections' containing title and content
        images: List of image dicts with 'id', 'data'
        output_file: Output PPTX path
        theme_name: Theme name (ocean, forest, sunset, purple, slate, coral) or "custom"
        title: Presentation title
        custom_theme: Custom theme dict (overrides theme_name if provided)
    """
    # Use custom theme if provided, otherwise use built-in
    if custom_theme:
        theme = custom_theme
    else:
        theme = THEMES.get(theme_name, THEMES["ocean"])

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print(f"Creating presentation with '{theme['name']}' theme...")

    # Title slide
    create_title_slide(prs, title, "AI-Generated Summary", theme)

    # Content slides
    sections = summary_data.get("sections", [])
    image_dict = {img["id"]: img["data"] for img in images}

    for idx, section in enumerate(sections, 1):
        print(f"  Creating slide {idx}/{len(sections)}: {section.get('title', 'Untitled')[:50]}...")

        slide_title = section.get("title", f"Section {idx}")
        content = section.get("content", "")
        linked_image_id = section.get("image_id")

        # Get image data if linked
        image_data = image_dict.get(linked_image_id) if linked_image_id else None

        create_content_slide(prs, slide_title, content, image_data, theme)

    # Save
    prs.save(output_file)
    print(f"\n✓ Presentation saved: {output_file}")
    print(f"  Theme: {theme['name']}")
    print(f"  Slides: {len(sections) + 1} (1 title + {len(sections)} content)")


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Generate modern PowerPoint presentations")
    parser.add_argument("summary_json", help="JSON file with summary sections")
    parser.add_argument("--output", default="presentation.pptx", help="Output file")
    parser.add_argument("--theme", default="ocean", choices=list(THEMES.keys()), help="Color theme")
    parser.add_argument("--title", default="Document Summary", help="Presentation title")

    args = parser.parse_args()

    # Load data
    with open(args.summary_json, 'r') as f:
        data = json.load(f)

    summary_data = data.get("summary_data", {})
    images = data.get("images", [])

    create_modern_presentation(summary_data, images, args.output, args.theme, args.title)
