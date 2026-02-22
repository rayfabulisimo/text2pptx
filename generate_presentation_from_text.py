#!/usr/bin/env python3
"""
Generate PowerPoint from skills_slides_text.md
"""

import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# Professional color theme
THEME = {
    "name": "Scientific Blue",
    "primary": RGBColor(13, 71, 161),      # Deep blue
    "secondary": RGBColor(25, 118, 210),   # Medium blue
    "accent": RGBColor(66, 165, 245),      # Light blue
    "background": RGBColor(250, 250, 252), # Very light blue-gray
    "text": RGBColor(33, 33, 33),          # Dark gray
    "text_on_dark": RGBColor(255, 255, 255)
}


def parse_markdown_slides(md_file):
    """Parse markdown file and extract slides."""
    with open(md_file, 'r') as f:
        content = f.read()

    # Split by slide markers (##)
    slides = []
    current_slide = None

    lines = content.split('\n')
    for line in lines:
        # Skip title and notes sections
        if line.startswith('# Pan-Cancer') or line.startswith('## Slide Text') or line.startswith('## Notes for'):
            continue

        # New slide
        if line.startswith('## Slide'):
            if current_slide:
                slides.append(current_slide)
            current_slide = {'title': '', 'bullets': []}

        # Title in bold
        elif line.startswith('**') and line.endswith('**') and current_slide is not None:
            current_slide['title'] = line.strip('*')

        # Bullet point
        elif line.startswith('- ') and current_slide is not None:
            current_slide['bullets'].append(line[2:])

        # Separator
        elif line.strip() == '---' and current_slide:
            slides.append(current_slide)
            current_slide = None

    # Add last slide if exists
    if current_slide:
        slides.append(current_slide)

    return slides


def create_title_slide(prs, title, subtitle):
    """Create title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Background
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 135
    fill.gradient_stops[0].color.rgb = THEME['background']
    fill.gradient_stops[1].color.rgb = RGBColor(255, 255, 255)

    # Large accent bar
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.5), Inches(7.5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = THEME['accent']
    accent.line.fill.background()

    # Title box
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(8.5), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = THEME['primary']
    p.alignment = PP_ALIGN.LEFT

    # Subtitle box
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.2),
        Inches(8.5), Inches(1)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.name = "Calibri"
    p.font.color.rgb = THEME['text']
    p.alignment = PP_ALIGN.LEFT


def create_content_slide(prs, title, bullets):
    """Create content slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = THEME['background']

    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.33), Inches(0.9)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = THEME['primary']
    header.line.fill.background()

    # Accent line
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.9),
        Inches(13.33), Inches(0.05)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = THEME['accent']
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2),
        Inches(12), Inches(0.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = THEME['text_on_dark']

    # Content area
    content_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.8),
        Inches(11.5), Inches(5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    # Add bullets
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
        p.font.name = "Calibri"
        p.font.color.rgb = THEME['text']
        p.space_before = Pt(12)
        p.space_after = Pt(6)
        p.line_spacing = 1.2


def create_presentation(md_file, output_file):
    """Create presentation from markdown file."""
    print(f"\nParsing {md_file}...")
    slides = parse_markdown_slides(md_file)
    print(f"Found {len(slides)} slides")

    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9
    prs.slide_height = Inches(7.5)

    print(f"\nCreating presentation with {THEME['name']} theme...")

    # First slide is title
    if slides and 'Title Slide' in str(slides[0]):
        create_title_slide(
            prs,
            "Pan-Cancer Single-Cell RNA-Seq",
            "Identifying Recurring Programs of Cellular Heterogeneity"
        )
        slides = slides[1:]  # Skip title slide

    # Content slides
    for idx, slide_data in enumerate(slides, 1):
        title = slide_data.get('title', f'Slide {idx}')
        bullets = slide_data.get('bullets', [])

        if bullets:  # Only create slide if it has content
            print(f"  Slide {idx+1}: {title[:60]}...")
            create_content_slide(prs, title, bullets)

    prs.save(output_file)
    print(f"\n✓ Presentation saved: {output_file}")
    print(f"  Theme: {THEME['name']}")
    print(f"  Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    md_file = "skills_slides_text.md"
    output_file = "pan_cancer_heterogeneity_presentation.pptx"

    create_presentation(md_file, output_file)
