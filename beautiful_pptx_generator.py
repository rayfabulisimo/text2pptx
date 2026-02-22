#!/usr/bin/env python3
"""
Beautiful PPTX Generator with Stunning Visuals
- Concise text (max 6 words per line)
- Beautiful backgrounds and icons
- Professional color schemes
"""

import sys
import io
from pathlib import Path
from typing import Dict, List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFilter

# Beautiful modern color schemes
BEAUTIFUL_THEMES = {
    "ocean": {
        "name": "Ocean Depths",
        "primary": RGBColor(0, 48, 73),        # Deep ocean blue
        "secondary": RGBColor(0, 119, 182),    # Ocean blue
        "accent": RGBColor(247, 127, 0),       # Coral orange
        "light": RGBColor(234, 246, 251),      # Light blue
        "text_on_dark": RGBColor(255, 255, 255),
        "text_on_light": RGBColor(33, 33, 33)
    },
    "forest": {
        "name": "Forest Canopy",
        "primary": RGBColor(20, 68, 56),       # Deep forest green
        "secondary": RGBColor(46, 125, 50),    # Fresh green
        "accent": RGBColor(255, 215, 0),       # Golden yellow
        "light": RGBColor(246, 251, 246),      # Pale green
        "text_on_dark": RGBColor(255, 255, 255),
        "text_on_light": RGBColor(33, 33, 33)
    },
    "sunset": {
        "name": "Sunset Glow",
        "primary": RGBColor(136, 14, 79),      # Deep magenta
        "secondary": RGBColor(216, 27, 96),    # Pink
        "accent": RGBColor(255, 193, 7),       # Amber
        "light": RGBColor(255, 241, 242),      # Pale pink
        "text_on_dark": RGBColor(255, 255, 255),
        "text_on_light": RGBColor(33, 33, 33)
    },
    "modern": {
        "name": "Modern Minimal",
        "primary": RGBColor(13, 27, 42),       # Nearly black
        "secondary": RGBColor(27, 38, 59),     # Dark blue-gray
        "accent": RGBColor(65, 105, 225),      # Royal blue
        "light": RGBColor(248, 249, 250),      # Off white
        "text_on_dark": RGBColor(255, 255, 255),
        "text_on_light": RGBColor(33, 33, 33)
    }
}


def create_beautiful_background(width=1200, height=900, theme_colors=None):
    """Create a beautiful gradient background with geometric patterns."""
    if theme_colors is None:
        theme_colors = {
            'primary': (0, 48, 73),
            'secondary': (0, 119, 182),
            'accent': (247, 127, 0)
        }

    img = Image.new('RGB', (width, height), theme_colors['primary'])
    draw = ImageDraw.Draw(img)

    # Create smooth gradient
    for i in range(height):
        ratio = i / height
        r = int(theme_colors['primary'][0] * (1-ratio) + theme_colors['secondary'][0] * ratio)
        g = int(theme_colors['primary'][1] * (1-ratio) + theme_colors['secondary'][1] * ratio)
        b = int(theme_colors['primary'][2] * (1-ratio) + theme_colors['secondary'][2] * ratio)
        draw.line([(0, i), (width, i)], fill=(r, g, b))

    # Add subtle geometric pattern
    for x in range(0, width, 100):
        for y in range(0, height, 100):
            # Semi-transparent circles
            draw.ellipse(
                [(x-20, y-20), (x+20, y+20)],
                fill=(255, 255, 255, 5)
            )

    # Apply slight blur for smoothness
    img = img.filter(ImageFilter.GaussianBlur(radius=2))

    return img


def create_icon(icon_type="circle", size=400, theme_colors=None):
    """Create simple, beautiful icons."""
    if theme_colors is None:
        theme_colors = {'accent': (247, 127, 0), 'light': (234, 246, 251)}

    img = Image.new('RGBA', (size, size), (255, 255, 255, 0))
    draw = ImageDraw.Draw(img)

    center = size // 2

    if icon_type == "circle":
        # Concentric circles
        for i in range(3):
            radius = center - (i * 50)
            color = theme_colors['accent'] if i % 2 == 0 else theme_colors['light']
            draw.ellipse(
                [(center-radius, center-radius), (center+radius, center+radius)],
                outline=color + (200,), width=8
            )

    elif icon_type == "star":
        # Star burst
        import math
        num_rays = 12
        for i in range(num_rays):
            angle = (2 * math.pi * i) / num_rays
            x = center + (center-20) * math.cos(angle)
            y = center + (center-20) * math.sin(angle)
            draw.line([(center, center), (x, y)], fill=theme_colors['accent'] + (200,), width=6)

    elif icon_type == "diamond":
        # Diamond shape
        points = [
            (center, 20),
            (size-20, center),
            (center, size-20),
            (20, center)
        ]
        draw.polygon(points, outline=theme_colors['accent'] + (200,), width=8)

    return img


def add_beautiful_background(slide, theme, bg_image=None):
    """Add beautiful background to slide."""
    if bg_image:
        # Use provided background image
        left = Inches(0)
        top = Inches(0)
        pic = slide.shapes.add_picture(
            bg_image, left, top,
            width=Inches(10), height=Inches(7.5)
        )
        # Send to back
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)
    else:
        # Solid color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = theme['light']


def create_title_slide_beautiful(prs, title, subtitle, theme, bg_image=None):
    """Create stunning title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    add_beautiful_background(slide, theme, bg_image)

    # Large color block
    block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(2.5),
        Inches(8), Inches(2.5)
    )
    block.fill.solid()
    block.fill.fore_color.rgb = theme['primary']
    block.line.fill.background()
    block.shadow.inherit = False

    # Accent bar
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(2.5),
        Inches(0.3), Inches(2.5)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = theme['accent']
    accent_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.8),
        Inches(7), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.name = "Arial"
    p.font.color.rgb = theme['text_on_dark']

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(4.2),
        Inches(7), Inches(0.6)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(16)
    p.font.name = "Arial"
    p.font.color.rgb = theme['text_on_dark']


def create_content_slide_beautiful(prs, title, bullets, theme, icon_type="circle", bg_image=None):
    """Create beautiful content slide with icon."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    add_beautiful_background(slide, theme, bg_image)

    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = theme['primary']
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25),
        Inches(9), Inches(0.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = "Arial"
    p.font.color.rgb = theme['text_on_dark']

    # Content area with bullets (left side)
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5),
        Inches(5.5), Inches(5.5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    for bullet in bullets:
        if tf.paragraphs[0].text:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
        p.font.name = "Arial"
        p.font.color.rgb = theme['text_on_light']
        p.space_before = Pt(12)

    # Add decorative icon on right side
    theme_colors = {
        'accent': tuple(theme['accent']),
        'light': tuple(theme['light'])
    }
    icon_img = create_icon(icon_type, 400, theme_colors)

    # Save to bytes
    icon_bytes = io.BytesIO()
    icon_img.save(icon_bytes, format='PNG')
    icon_bytes.seek(0)

    # Add to slide
    icon_pic = slide.shapes.add_picture(
        icon_bytes,
        Inches(6.5), Inches(2.5),
        width=Inches(3)
    )


def create_beautiful_presentation(
    title: str,
    sections: List[Dict],
    output_file: str,
    theme_name: str = "ocean"
):
    """
    Create beautiful presentation with concise text.

    sections: List of dicts with 'title' and 'bullets' (list of strings, max 6 words each)
    """
    theme = BEAUTIFUL_THEMES.get(theme_name, BEAUTIFUL_THEMES["ocean"])

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print(f"\n🎨 Creating presentation with '{theme['name']}' theme...")

    # Create background image
    theme_colors = {
        'primary': tuple(theme['primary']),
        'secondary': tuple(theme['secondary']),
        'accent': tuple(theme['accent'])
    }
    bg_img = create_beautiful_background(1200, 900, theme_colors)
    bg_bytes = io.BytesIO()
    bg_img.save(bg_bytes, format='PNG')
    bg_bytes.seek(0)

    # Title slide
    create_title_slide_beautiful(prs, title, "Research Summary", theme, bg_bytes)

    # Content slides
    icon_types = ["circle", "star", "diamond"]
    for idx, section in enumerate(sections):
        icon_type = icon_types[idx % len(icon_types)]
        print(f"   Slide {idx+2}: {section['title'][:50]}...")

        bg_bytes.seek(0)
        create_content_slide_beautiful(
            prs,
            section['title'],
            section['bullets'],
            theme,
            icon_type,
            bg_bytes
        )

    prs.save(output_file)
    print(f"\n✓ Presentation saved: {output_file}")
    print(f"  Theme: {theme['name']}")
    print(f"  Slides: {len(sections) + 1}")


if __name__ == "__main__":
    # Test
    sections = [
        {
            "title": "Key Finding",
            "bullets": [
                "Found twelve expression programs",
                "Cell lines mirror tumors",
                "Validates experimental models"
            ]
        }
    ]

    create_beautiful_presentation(
        "Test Presentation",
        sections,
        "test_beautiful.pptx",
        "ocean"
    )
