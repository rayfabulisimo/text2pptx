#!/usr/bin/env python3
"""
Document Analyzer Agent
Handles multiple document formats: PDF, DOCX, PPTX, TXT, MD, images
Extracts text, images, and creates structured summaries
"""

import asyncio
import sys
import os
import json
from pathlib import Path
from typing import List, Dict, Optional, Tuple
import io
import base64

# Document processing imports
try:
    import fitz  # PyMuPDF for PDFs
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    from docx import Document
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import _Cell, Table
    from docx.text.paragraph import Paragraph
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation as PptxPresentation
    PPTX_READ_AVAILABLE = True
except ImportError:
    PPTX_READ_AVAILABLE = False

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False


def extract_text_from_pdf(file_path: str, start_page: int = 1, end_page: int = None) -> Tuple[str, List[Dict]]:
    """Extract text and images from PDF."""
    if not PYMUPDF_AVAILABLE:
        return "Error: PyMuPDF not installed", []

    doc = fitz.open(file_path)
    text_content = []
    images = []

    end = end_page if end_page else len(doc)

    for page_num in range(start_page - 1, end):
        page = doc[page_num]

        # Extract text
        text_content.append(page.get_text())

        # Extract images
        image_list = page.get_images(full=True)
        for img_idx, img in enumerate(image_list):
            xref = img[0]
            base_image = doc.extract_image(xref)

            images.append({
                "source_page": page_num + 1,
                "index": img_idx,
                "id": f"pdf_p{page_num+1}_i{img_idx}",
                "format": base_image["ext"],
                "data": base_image["image"]
            })

    doc.close()
    return "\n\n".join(text_content), images


def extract_text_from_docx(file_path: str) -> Tuple[str, List[Dict]]:
    """Extract text and images from DOCX."""
    if not DOCX_AVAILABLE:
        return "Error: python-docx not installed", []

    doc = Document(file_path)
    text_content = []
    images = []

    # Extract text from paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            text_content.append(para.text)

    # Extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text for cell in row.cells]
            text_content.append(" | ".join(row_text))

    # Extract images
    for rel_id, rel in doc.part.rels.items():
        if "image" in rel.target_ref:
            try:
                image_data = rel.target_part.blob
                images.append({
                    "source_page": "document",
                    "index": len(images),
                    "id": f"docx_i{len(images)}",
                    "format": "png",
                    "data": image_data
                })
            except:
                pass

    return "\n\n".join(text_content), images


def extract_text_from_pptx(file_path: str) -> Tuple[str, List[Dict]]:
    """Extract text and images from PPTX."""
    if not PPTX_READ_AVAILABLE:
        return "Error: python-pptx not installed", []

    prs = PptxPresentation(file_path)
    text_content = []
    images = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []

        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)

            # Extract images
            if shape.shape_type == 13:  # Picture
                try:
                    image = shape.image
                    images.append({
                        "source_page": slide_num,
                        "index": len(images),
                        "id": f"pptx_s{slide_num}_i{len(images)}",
                        "format": image.ext,
                        "data": image.blob
                    })
                except:
                    pass

        if slide_text:
            text_content.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))

    return "\n\n".join(text_content), images


def extract_text_from_txt_or_md(file_path: str) -> Tuple[str, List[Dict]]:
    """Extract text from plain text or markdown files."""
    with open(file_path, 'r', encoding='utf-8') as f:
        text = f.read()
    return text, []


def extract_text_from_image(file_path: str) -> Tuple[str, List[Dict]]:
    """Handle image file - store as image, no text extraction."""
    if not PIL_AVAILABLE:
        return "Error: Pillow not installed", []

    with open(file_path, 'rb') as f:
        image_data = f.read()

    images = [{
        "source_page": "image_file",
        "index": 0,
        "id": "img_0",
        "format": Path(file_path).suffix[1:],
        "data": image_data
    }]

    return f"[Image file: {Path(file_path).name}]", images


def analyze_document(file_path: str, pages: str = None) -> Dict:
    """
    Analyze document and extract content.

    Args:
        file_path: Path to document
        pages: Page range (PDF only, e.g., "1-27")

    Returns:
        Dict with text, images, and metadata
    """
    file_path = Path(file_path)

    if not file_path.exists():
        return {"error": f"File not found: {file_path}"}

    ext = file_path.suffix.lower()

    # Parse page range for PDFs
    start_page, end_page = 1, None
    if pages and ext == ".pdf":
        try:
            parts = pages.split("-")
            start_page = int(parts[0])
            end_page = int(parts[1]) if len(parts) > 1 else None
        except:
            pass

    # Extract based on file type
    if ext == ".pdf":
        text, images = extract_text_from_pdf(str(file_path), start_page, end_page)
    elif ext in [".docx", ".doc"]:
        text, images = extract_text_from_docx(str(file_path))
    elif ext in [".pptx", ".ppt"]:
        text, images = extract_text_from_pptx(str(file_path))
    elif ext in [".txt", ".md"]:
        text, images = extract_text_from_txt_or_md(str(file_path))
    elif ext in [".png", ".jpg", ".jpeg", ".gif", ".bmp"]:
        text, images = extract_text_from_image(str(file_path))
    else:
        return {"error": f"Unsupported file format: {ext}"}

    return {
        "file_path": str(file_path),
        "file_name": file_path.name,
        "file_type": ext[1:],
        "text": text,
        "images": images,
        "num_images": len(images),
        "text_length": len(text)
    }


async def main():
    """Main entry point."""
    import argparse

    parser = argparse.ArgumentParser(
        description="Analyze documents and extract text and images",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Supported formats:
  - PDF (.pdf)
  - Word (.docx, .doc)
  - PowerPoint (.pptx, .ppt)
  - Text (.txt, .md)
  - Images (.png, .jpg, .jpeg, .gif, .bmp)

Examples:
  python document_analyzer_agent.py paper.pdf --pages 1-27
  python document_analyzer_agent.py document.docx
  python document_analyzer_agent.py presentation.pptx
  python document_analyzer_agent.py report.txt
        """
    )
    parser.add_argument("file_path", help="Path to document file")
    parser.add_argument("--pages", help="Page range for PDFs (e.g., 1-27)", default=None)
    parser.add_argument("--output", help="Output JSON file", default="document_analysis.json")

    args = parser.parse_args()

    print(f"Analyzing: {args.file_path}")
    print(f"File type: {Path(args.file_path).suffix}")

    # Analyze document
    result = analyze_document(args.file_path, args.pages)

    if "error" in result:
        print(f"\n❌ Error: {result['error']}")
        sys.exit(1)

    # Display results
    print(f"\n✓ Analysis complete!")
    print(f"  Text extracted: {result['text_length']:,} characters")
    print(f"  Images found: {result['num_images']}")

    # Save results (without image data for readability)
    output_data = {
        "file_path": result["file_path"],
        "file_name": result["file_name"],
        "file_type": result["file_type"],
        "text": result["text"],
        "text_length": result["text_length"],
        "num_images": result["num_images"],
        "images": [
            {k: v for k, v in img.items() if k != "data"}
            for img in result["images"]
        ]
    }

    with open(args.output, "w") as f:
        json.dump(output_data, f, indent=2)

    print(f"\n✓ Analysis saved to: {args.output}")
    print(f"\n  Preview of text (first 500 chars):")
    print(f"  {result['text'][:500]}...")


if __name__ == "__main__":
    asyncio.run(main())
